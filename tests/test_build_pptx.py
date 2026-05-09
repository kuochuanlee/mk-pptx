"""
build_pptx.py 整合測試。

測試策略：
  - 使用現有的 template 檔案（template/Creative-Idea-Bulb-PowerPoint-Template.pptx）
  - 每種 layout_type 各一個測試，驗證 .pptx 可正常開啟、slide 數量正確、title 正確
  - diagram 相關測試：mock build_diagram，不依賴真實 mmdc
  - 所有測試的產出 .pptx 放在 tmp_path 暫存目錄

CI 注意事項：
  - 這些測試依賴 template 檔案，若 template 不存在則 skip
  - diagram 相關測試透過 mock 隔離 mmdc 依賴
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PptxPresentation

from schemas.slide_schema import Presentation, Slide, SlideContent, LayoutType
from src.build_pptx import build_pptx
from src.build_diagram import DiagramResult

# ============================================================
# 常數與 fixtures
# ============================================================

# 模板與設定檔路徑（相對於專案根目錄）
TEMPLATE_PATH = Path("template") / "Creative-Idea-Bulb-PowerPoint-Template.pptx"
TEMPLATE_MAP_PATH = Path("config") / "template_map.json"

# 所有依賴模板的測試共用此 skip 條件
requires_template = pytest.mark.skipif(
    not TEMPLATE_PATH.exists(),
    reason="模板檔案不存在，跳過此測試",
)


def _make_presentation(*slides: Slide) -> Presentation:
    """建立 Presentation 測試物件。"""
    return Presentation(
        presentation_title="測試簡報",
        author="測試人員",
        date="2025-06-01",
        slides=list(slides),
    )


def _make_slide(
    slide_number: int,
    layout_type: str,
    title: str,
    content: dict,
    speaker_notes: str | None = None,
) -> Slide:
    """建立 Slide 測試物件。"""
    data = {
        "slide_number": slide_number,
        "layout_type": layout_type,
        "title": title,
        "content": content,
    }

    if speaker_notes is not None:
        data["speaker_notes"] = speaker_notes

    return Slide.model_validate(data)


# ============================================================
# 輔助函式：讀取產出的 .pptx 並驗證基本結構
# ============================================================


def _open_pptx(path: Path) -> PptxPresentation:
    """開啟 .pptx 並回傳 Presentation 物件。若無法開啟則讓測試失敗。"""
    assert path.exists(), f"輸出的 .pptx 不存在：{path}"
    return PptxPresentation(str(path))


def _get_slide_title(prs: PptxPresentation, slide_index: int) -> str:
    """取得指定投影片的標題文字（idx=10 placeholder）。"""
    slide = prs.slides[slide_index]

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            return ph.text

    return ""


# ============================================================
# title_slide 測試
# ============================================================


@requires_template
class TestBuildPptxTitleSlide:
    """title_slide layout_type 的 PPTX 生成測試。"""

    def test_generates_pptx_file(self, tmp_path: Path):
        """應產出 .pptx 檔案。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title="Q2 營運報告",
            content={"subtitle": "2025 年第二季"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        result_path = build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        assert result_path.exists()

    def test_slide_count_is_correct(self, tmp_path: Path):
        """產出的 .pptx 投影片數量應等於 slides 清單長度。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title="Q2 營運報告",
            content={"subtitle": "2025 年第二季"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        # 模板範例 slides 已被清除，產出的 slide 數量應等於新增的數量
        prs = _open_pptx(output)

        assert len(prs.slides) == 1

    def test_title_text_is_correct(self, tmp_path: Path):
        """產出的 .pptx 最後一頁（新增的 slide）標題應正確。"""
        title_text = "Q2 營運報告"
        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title=title_text,
            content={"subtitle": "2025 年第二季"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        # 新增的 slide 在最後一頁
        last_idx = len(prs.slides) - 1
        actual_title = _get_slide_title(prs, last_idx)

        assert actual_title == title_text


# ============================================================
# section_divider 測試
# ============================================================


@requires_template
class TestBuildPptxSectionDivider:
    """section_divider layout_type 的 PPTX 生成測試。"""

    def test_generates_without_subtitle(self, tmp_path: Path):
        """section_divider 不含 subtitle 時應正常生成。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="section_divider",
            title="第一章：市場分析",
            content={},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_generates_with_subtitle(self, tmp_path: Path):
        """section_divider 含 subtitle 時應正常生成。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="section_divider",
            title="第一章",
            content={"subtitle": "市場分析概述"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0


# ============================================================
# bullets_only 測試
# ============================================================


@requires_template
class TestBuildPptxBulletsOnly:
    """bullets_only layout_type 的 PPTX 生成測試。"""

    def test_generates_with_bullets(self, tmp_path: Path):
        """應正常生成含條列的 .pptx。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_only",
            title="重點摘要",
            content={"bullets": ["整體成長 15%", "三個區域均達標", "成本下降 7%"]},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_title_is_correct(self, tmp_path: Path):
        """bullets_only 投影片標題應正確。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_only",
            title="重點摘要",
            content={"bullets": ["條列一"]},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        last_idx = len(prs.slides) - 1
        assert _get_slide_title(prs, last_idx) == "重點摘要"

    def test_speaker_notes_is_filled(self, tmp_path: Path):
        """speaker_notes 有值時應填入投影片備註。"""
        notes_text = "請特別強調成本下降的原因。"
        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_only",
            title="重點",
            content={"bullets": ["條列一"]},
            speaker_notes=notes_text,
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        last_slide = prs.slides[-1]
        actual_notes = last_slide.notes_slide.notes_text_frame.text

        assert actual_notes == notes_text


# ============================================================
# table_only 測試
# ============================================================


@requires_template
class TestBuildPptxTableOnly:
    """table_only layout_type 的 PPTX 生成測試。"""

    def test_generates_with_table(self, tmp_path: Path):
        """應正常生成含表格的 .pptx。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="table_only",
            title="數據比較",
            content={
                "markdown_table": (
                    "| 項目 | Q1 | Q2 |\n"
                    "|---|---|---|\n"
                    "| 營收 | 100 | 115 |\n"
                    "| 成本 | 80 | 74 |"
                )
            },
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_table_shape_exists_in_slide(self, tmp_path: Path):
        """產出的投影片應包含一個表格 shape。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="table_only",
            title="數據比較",
            content={
                "markdown_table": "| A | B |\n|---|---|\n| 1 | 2 |"
            },
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        last_slide = prs.slides[-1]

        # 檢查是否有 table shape
        has_table = any(shape.has_table for shape in last_slide.shapes)
        assert has_table


# ============================================================
# bullets_with_table 測試
# ============================================================


@requires_template
class TestBuildPptxBulletsWithTable:
    """bullets_with_table layout_type 的 PPTX 生成測試。"""

    def test_generates_with_bullets_and_table(self, tmp_path: Path):
        """應正常生成含條列和表格的 .pptx。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_with_table",
            title="綜合分析",
            content={
                "bullets": ["整體成長 15%", "三個區域均達標"],
                "markdown_table": (
                    "| 項目 | Q1 | Q2 |\n"
                    "|---|---|---|\n"
                    "| 營收 | 100 | 115 |"
                ),
            },
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        last_slide = prs.slides[-1]

        # 應同時有 textbox 和 table
        has_table = any(shape.has_table for shape in last_slide.shapes)
        assert has_table

    def test_slide_count_multiple_slides(self, tmp_path: Path):
        """多頁簡報的 slide 總數應正確。"""
        slides = [
            _make_slide(
                slide_number=1,
                layout_type="title_slide",
                title="封面",
                content={"subtitle": "副標題"},
            ),
            _make_slide(
                slide_number=2,
                layout_type="bullets_with_table",
                title="分析",
                content={
                    "bullets": ["重點一"],
                    "markdown_table": "| A |\n|---|\n| 1 |",
                },
            ),
        ]
        prs_data = _make_presentation(*slides)
        output = tmp_path / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)

        # 模板範例 slides 已被清除，只有新增的 2 頁
        assert len(prs.slides) == 2


# ============================================================
# diagram_only 測試（mock build_diagram，不需要真實 mmdc）
# ============================================================


@requires_template
class TestBuildPptxDiagramOnly:
    """diagram_only layout_type 的 PPTX 生成測試，mock build_diagram。"""

    def _mock_build_diagram_success(self, tmp_path: Path):
        """建立模擬成功的 build_diagram mock，回傳真實存在的假 PNG。"""
        # 建立一個最小的有效 PNG 檔（1x1 透明 PNG）
        fake_png = tmp_path / "fake_diagram.png"
        fake_png.write_bytes(
            b"\x89PNG\r\n\x1a\n"  # PNG signature
            b"\x00\x00\x00\rIHDR"  # IHDR chunk
            b"\x00\x00\x00\x01"   # width=1
            b"\x00\x00\x00\x01"   # height=1
            b"\x08\x02"           # 8-bit depth, RGB
            b"\x00\x00\x00"       # compression, filter, interlace
            b"\x90wS\xde"         # CRC
            b"\x00\x00\x00\x0cIDATx"  # IDAT chunk
            b"\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
            b"\x00\x00\x00\x00IEND\xaeB`\x82"  # IEND chunk
        )

        mock_result = MagicMock(spec=DiagramResult)
        mock_result.png_path = fake_png
        mock_result.svg_path = tmp_path / "fake_diagram.svg"

        return mock_result

    def test_generates_with_diagram_success(self, tmp_path: Path):
        """mmdc 渲染成功時應正常生成 .pptx。"""
        mock_result = self._mock_build_diagram_success(tmp_path)

        slide = _make_slide(
            slide_number=1,
            layout_type="diagram_only",
            title="流程圖",
            content={"mermaid": "flowchart LR\n  A --> B"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        with patch("src.build_pptx.build_diagram", return_value=mock_result):
            build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_fallback_when_diagram_fails(self, tmp_path: Path):
        """mmdc 渲染失敗（回傳 None）時應 fallback，不拋例外。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="diagram_only",
            title="流程圖",
            content={"mermaid": "flowchart LR\n  A --> B"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        # build_diagram 回傳 None 模擬渲染失敗
        with patch("src.build_pptx.build_diagram", return_value=None):
            build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        # 應成功生成 .pptx（fallback 為文字框）
        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_title_is_correct_when_diagram_fails(self, tmp_path: Path):
        """fallback 時投影片標題應仍正確。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="diagram_only",
            title="流程圖",
            content={"mermaid": "flowchart LR\n  A --> B"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        with patch("src.build_pptx.build_diagram", return_value=None):
            build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        last_idx = len(prs.slides) - 1
        assert _get_slide_title(prs, last_idx) == "流程圖"


# ============================================================
# bullets_with_diagram 測試（mock build_diagram）
# ============================================================


@requires_template
class TestBuildPptxBulletsWithDiagram:
    """bullets_with_diagram layout_type 的 PPTX 生成測試，mock build_diagram。"""

    def test_generates_with_diagram_success(self, tmp_path: Path):
        """mmdc 渲染成功時應正常生成包含條列和圖表的 .pptx。"""
        # 建立假 PNG
        fake_png = tmp_path / "fake.png"
        fake_png.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
            b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02"
            b"\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
            b"\x00\x00\x00\x00IEND\xaeB`\x82"
        )

        mock_result = MagicMock(spec=DiagramResult)
        mock_result.png_path = fake_png

        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_with_diagram",
            title="流程說明",
            content={
                "bullets": ["步驟一：送出申請", "步驟二：系統驗證"],
                "mermaid": "flowchart LR\n  A --> B --> C",
            },
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        with patch("src.build_pptx.build_diagram", return_value=mock_result):
            build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0

    def test_fallback_when_diagram_fails(self, tmp_path: Path):
        """mmdc 渲染失敗時應 fallback 顯示文字，不拋例外。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="bullets_with_diagram",
            title="流程說明",
            content={
                "bullets": ["步驟一"],
                "mermaid": "flowchart LR\n  A --> B",
            },
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        with patch("src.build_pptx.build_diagram", return_value=None):
            build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        prs = _open_pptx(output)
        assert len(prs.slides) > 0


# ============================================================
# 錯誤處理測試
# ============================================================


class TestBuildPptxErrors:
    """build_pptx 錯誤情境測試。"""

    def test_raises_when_template_not_found(self, tmp_path: Path):
        """模板檔案不存在時應拋出 FileNotFoundError。"""
        # 建立一個指向不存在模板的 template_map
        fake_map = tmp_path / "template_map.json"
        fake_map.write_text(
            '{"template_path": "nonexistent.pptx", '
            '"slide_width_emu": 9144000, "slide_height_emu": 5143500, '
            '"layouts": {"title_slide": {"layout_name": "Cover Slide layout", '
            '"title_idx": 10, "subtitle_idx": 11}}}'
        )

        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title="測試",
            content={"subtitle": "副標題"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"

        with pytest.raises(FileNotFoundError):
            build_pptx(prs_data, output, fake_map)

    def test_raises_when_template_map_not_found(self, tmp_path: Path):
        """template_map.json 不存在時應拋出 FileNotFoundError。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title="測試",
            content={"subtitle": "副標題"},
        )
        prs_data = _make_presentation(slide)
        output = tmp_path / "out.pptx"
        nonexistent_map = tmp_path / "nonexistent_map.json"

        with pytest.raises(FileNotFoundError):
            build_pptx(prs_data, output, nonexistent_map)

    @requires_template
    def test_output_path_parent_created_automatically(self, tmp_path: Path):
        """輸出路徑的上層目錄不存在時應自動建立。"""
        slide = _make_slide(
            slide_number=1,
            layout_type="title_slide",
            title="測試",
            content={"subtitle": "副標題"},
        )
        prs_data = _make_presentation(slide)
        # 使用不存在的子目錄
        output = tmp_path / "nested" / "deep" / "out.pptx"

        build_pptx(prs_data, output, TEMPLATE_MAP_PATH)

        assert output.exists()
