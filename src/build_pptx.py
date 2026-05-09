"""
JSON -> .pptx 核心邏輯。

讀取 template_map.json 取得模板路徑與 placeholder 對應，
依據 Presentation 物件中每頁 Slide 的 layout_type
分派至對應的處理函式，組裝完整的 .pptx 檔案。

各 layout_type 處理方式：
  - title_slide       : 填 title placeholder + subtitle placeholder
  - section_divider   : 填 title placeholder（+ 可選 subtitle）
  - bullets_only      : title placeholder + add_textbox() 條列
  - table_only        : title placeholder + add_table() 表格
  - bullets_with_table: title + 左側條列 textbox + 右側表格
  - diagram_only      : title + add_picture() PNG（失敗 fallback 條列）
  - bullets_with_diagram: title + 左側條列 + 右側 PNG（失敗 fallback 條列）
"""

from __future__ import annotations

import json
import logging
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from schemas.slide_schema import LayoutType, Presentation, Slide
from src.build_diagram import build_diagram
from src.table_parser import parse_markdown_table

logger = logging.getLogger(__name__)

# 公司主色（標題列背景）
COLOR_PRIMARY = RGBColor(0x00, 0x5B, 0xAC)

# 表格資料列底色
COLOR_ROW_EVEN = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ODD = RGBColor(0xF2, 0xF2, 0xF2)

# 表格文字白色
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 表格文字黑色
COLOR_BLACK = RGBColor(0x33, 0x33, 0x33)

# 字型名稱
FONT_BODY = "Microsoft JhengHei"

# 字級
FONT_SIZE_TITLE_ROW = Pt(14)
FONT_SIZE_BODY = Pt(12)

# 預設 template_map.json 路徑（相對於專案根目錄）
DEFAULT_TEMPLATE_MAP_PATH = Path("config") / "template_map.json"

# 圖表輸出目錄
DIAGRAM_OUTPUT_DIR = Path("output")

# 圖表 mmdc 設定檔路徑
DIAGRAM_CONFIG_PATH = Path("config") / "mermaid.json"

# 版面邊距（EMU）
MARGIN_LEFT = Emu(457200)   # 0.5 inch
MARGIN_TOP_CONTENT = Emu(1143000)  # 1.25 inch（標題下方）
MARGIN_RIGHT = Emu(457200)  # 0.5 inch

# bullet 符號
BULLET_CHAR = "\u2022"


# ============================================================
# 主要入口
# ============================================================


def build_pptx(
    presentation_data: Presentation,
    output_path: Path,
    template_map_path: Path | None = None,
) -> Path:
    """將驗證過的 Presentation 物件轉換為 .pptx 檔案。

    Args:
        presentation_data: 已通過 Pydantic 驗證的 Presentation 物件
        output_path: 輸出 .pptx 檔案路徑
        template_map_path: template_map.json 路徑（None 使用預設）

    Returns:
        實際輸出的 .pptx 檔案 Path

    Raises:
        FileNotFoundError: 模板檔案或 template_map.json 不存在
        ValueError: template_map.json 中找不到對應的 layout_name
    """
    # 讀取 template_map.json
    if template_map_path is None:
        template_map_path = DEFAULT_TEMPLATE_MAP_PATH

    template_map = _load_template_map(template_map_path)

    # 載入 PPTX 模板
    template_path = Path(template_map["template_path"])

    if not template_path.exists():
        raise FileNotFoundError(f"模板檔案不存在：{template_path}")

    prs = PptxPresentation(str(template_path))

    # 取得投影片尺寸（供版面計算使用）
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 建立 layout_name -> slide layout 的對應表。
    # 部分模板的 slide_layouts 僅包含少數 layout，
    # 其餘 layout 只能從現有 slides 的 slide_layout 屬性取得。
    layout_map: dict = {}

    # 先從 prs.slide_layouts 收集（官方清單）
    for sl in prs.slide_layouts:
        layout_map[sl.name] = sl

    # 再從現有 slides 的 slide_layout 補充（含模板中嵌入但未列在 slide_layouts 的 layout）
    for existing_slide in prs.slides:
        sl = existing_slide.slide_layout
        if sl.name not in layout_map:
            layout_map[sl.name] = sl

    # 清除模板中的範例 slides（只保留 layout 定義）
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 依序處理每一頁 slide
    for slide_data in presentation_data.slides:
        _process_slide(
            prs=prs,
            slide_data=slide_data,
            template_map=template_map,
            layout_map=layout_map,
            slide_width=slide_width,
            slide_height=slide_height,
        )

    # 確保輸出目錄存在
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(output_path))

    logger.info(
        "PPTX 生成完成：%s（共 %d 頁）",
        output_path,
        len(presentation_data.slides),
    )

    return output_path


# ============================================================
# 設定檔載入
# ============================================================


def _load_template_map(template_map_path: Path) -> dict:
    """讀取並回傳 template_map.json 的內容。

    Args:
        template_map_path: template_map.json 的 Path

    Returns:
        dict 形式的設定內容

    Raises:
        FileNotFoundError: 設定檔不存在
    """
    if not template_map_path.exists():
        raise FileNotFoundError(f"template_map.json 不存在：{template_map_path}")

    with template_map_path.open(encoding="utf-8") as f:
        return json.load(f)


def _get_layout_config(template_map: dict, layout_type: LayoutType) -> dict:
    """從 template_map 取得特定 layout_type 的設定。

    Args:
        template_map: template_map.json 的 dict
        layout_type: LayoutType 列舉值

    Returns:
        該 layout_type 的設定 dict（含 layout_name, title_idx, subtitle_idx）

    Raises:
        KeyError: layout_type 不在 template_map 中
    """
    layouts = template_map.get("layouts", {})
    key = layout_type.value

    if key not in layouts:
        raise KeyError(
            f"template_map.json 中找不到 layout_type '{key}' 的設定。"
            f"可用的 layout_type：{list(layouts.keys())}"
        )

    return layouts[key]


def _find_slide_layout(layout_map: dict, layout_name: str):
    """在 layout_map 中尋找對應的 slide layout 物件。

    Args:
        layout_map: layout_name -> slide layout 的對應表
        layout_name: 要尋找的 layout 名稱

    Returns:
        pptx SlideLayout 物件

    Raises:
        ValueError: 找不到對應的 layout_name
    """
    if layout_name not in layout_map:
        raise ValueError(
            f"模板中找不到 layout '{layout_name}'。"
            f"可用的 layout：{list(layout_map.keys())}"
        )

    return layout_map[layout_name]


# ============================================================
# 投影片分派處理
# ============================================================


def _process_slide(
    prs: PptxPresentation,
    slide_data: Slide,
    template_map: dict,
    layout_map: dict,
    slide_width: int,
    slide_height: int,
) -> None:
    """依 layout_type 分派至對應的處理函式。

    Args:
        prs: python-pptx Presentation 物件
        slide_data: 單頁 Slide 資料
        template_map: template_map.json 設定
        layout_map: layout_name -> slide layout 對應表
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    # 取得此 layout_type 的設定
    layout_cfg = _get_layout_config(template_map, slide_data.layout_type)
    layout_name = layout_cfg["layout_name"]
    title_idx = layout_cfg["title_idx"]
    subtitle_idx = layout_cfg["subtitle_idx"]

    # 找到對應的 slide layout 並新增投影片
    slide_layout = _find_slide_layout(layout_map, layout_name)
    slide = prs.slides.add_slide(slide_layout)

    # 填入 title（所有 layout 共通）
    _fill_placeholder(slide, title_idx, slide_data.title)

    # 依 layout_type 分派處理
    layout = slide_data.layout_type

    if layout == LayoutType.TITLE_SLIDE:
        _handle_title_slide(slide, slide_data, subtitle_idx)

    elif layout == LayoutType.SECTION_DIVIDER:
        _handle_section_divider(slide, slide_data, subtitle_idx)

    elif layout == LayoutType.BULLETS_ONLY:
        _handle_bullets_only(slide, slide_data, slide_width, slide_height)

    elif layout == LayoutType.TABLE_ONLY:
        _handle_table_only(slide, slide_data, slide_width, slide_height)

    elif layout == LayoutType.BULLETS_WITH_TABLE:
        _handle_bullets_with_table(slide, slide_data, slide_width, slide_height)

    elif layout == LayoutType.DIAGRAM_ONLY:
        _handle_diagram_only(slide, slide_data, slide_width, slide_height)

    elif layout == LayoutType.BULLETS_WITH_DIAGRAM:
        _handle_bullets_with_diagram(slide, slide_data, slide_width, slide_height)

    # 填入 speaker_notes（若有）
    if slide_data.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes


# ============================================================
# Placeholder 填入
# ============================================================


def _fill_placeholder(slide, idx: int, text: str) -> None:
    """將文字填入指定 idx 的 placeholder。

    若 placeholder 不存在則記錄警告並略過，不中斷流程。

    Args:
        slide: pptx Slide 物件
        idx: placeholder idx
        text: 要填入的文字
    """
    try:
        ph = slide.placeholders[idx]
        ph.text = text
    except KeyError:
        logger.warning("slide placeholder idx=%d 不存在，略過填入。", idx)


# ============================================================
# 各 layout_type 處理函式
# ============================================================


def _handle_title_slide(slide, slide_data: Slide, subtitle_idx: int) -> None:
    """處理 title_slide：填入 subtitle placeholder。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        subtitle_idx: subtitle placeholder 的 idx
    """
    subtitle = slide_data.content.subtitle or ""
    _fill_placeholder(slide, subtitle_idx, subtitle)


def _handle_section_divider(slide, slide_data: Slide, subtitle_idx: int) -> None:
    """處理 section_divider：選擇性填入 subtitle placeholder。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        subtitle_idx: subtitle placeholder 的 idx
    """
    subtitle = slide_data.content.subtitle or ""

    if subtitle:
        _fill_placeholder(slide, subtitle_idx, subtitle)


def _handle_bullets_only(
    slide,
    slide_data: Slide,
    slide_width: int,
    slide_height: int,
) -> None:
    """處理 bullets_only：在標題下方建立條列文字框。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    bullets = slide_data.content.bullets or []

    # 計算文字框位置與尺寸
    left = MARGIN_LEFT
    top = MARGIN_TOP_CONTENT
    width = Emu(slide_width) - MARGIN_LEFT - MARGIN_RIGHT
    height = Emu(slide_height) - top - MARGIN_RIGHT

    _add_bullet_textbox(slide, bullets, left, top, width, height)


def _handle_table_only(
    slide,
    slide_data: Slide,
    slide_width: int,
    slide_height: int,
) -> None:
    """處理 table_only：解析 markdown_table 並建立 PPTX 表格。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    md_table = slide_data.content.markdown_table or ""
    headers, rows = parse_markdown_table(md_table)

    # 表格置中，寬度佔投影片 80%
    table_width = int(slide_width * 0.80)
    left = int((slide_width - table_width) / 2)
    top = int(MARGIN_TOP_CONTENT)

    # 計算列數：header + 資料列
    row_count = 1 + len(rows)
    col_count = len(headers)

    _add_styled_table(
        slide=slide,
        headers=headers,
        rows=rows,
        left=left,
        top=top,
        width=table_width,
        row_count=row_count,
        col_count=col_count,
        slide_height=slide_height,
    )


def _handle_bullets_with_table(
    slide,
    slide_data: Slide,
    slide_width: int,
    slide_height: int,
) -> None:
    """處理 bullets_with_table：左側條列 + 右側表格。

    左半區寬度 45%，右半區寬度 50%，中間 5% 為間距。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    bullets = slide_data.content.bullets or []
    md_table = slide_data.content.markdown_table or ""
    headers, rows = parse_markdown_table(md_table)

    top = int(MARGIN_TOP_CONTENT)
    content_height = slide_height - top - int(MARGIN_RIGHT)

    # 左側條列（45% 寬）
    left_width = int(slide_width * 0.45)
    left_left = int(MARGIN_LEFT)

    _add_bullet_textbox(
        slide=slide,
        bullets=bullets,
        left=left_left,
        top=top,
        width=left_width,
        height=content_height,
    )

    # 右側表格（50% 寬，從 50% 位置開始）
    right_left = int(slide_width * 0.50)
    right_width = int(slide_width * 0.50) - int(MARGIN_RIGHT)

    row_count = 1 + len(rows)
    col_count = len(headers)

    _add_styled_table(
        slide=slide,
        headers=headers,
        rows=rows,
        left=right_left,
        top=top,
        width=right_width,
        row_count=row_count,
        col_count=col_count,
        slide_height=slide_height,
    )


def _handle_diagram_only(
    slide,
    slide_data: Slide,
    slide_width: int,
    slide_height: int,
) -> None:
    """處理 diagram_only：渲染 mermaid 為 PNG 並插入投影片置中。

    若渲染失敗，fallback 為顯示 mermaid 原始碼的文字框。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    mermaid_string = slide_data.content.mermaid or ""

    # 嘗試渲染圖表
    output_stem = f"slide_{slide_data.slide_number}_diagram"
    result = build_diagram(
        mermaid_string=mermaid_string,
        output_stem=output_stem,
        output_dir=DIAGRAM_OUTPUT_DIR,
        config_path=DIAGRAM_CONFIG_PATH,
    )

    if result is not None:
        # 渲染成功：PNG 插入投影片，寬度佔 70%，垂直置中
        pic_width = int(slide_width * 0.70)
        pic_left = int((slide_width - pic_width) / 2)
        top = int(MARGIN_TOP_CONTENT)

        slide.shapes.add_picture(
            str(result.png_path),
            pic_left,
            top,
            width=pic_width,
        )

        return

    # 渲染失敗：fallback 顯示 mermaid 原始碼
    logger.warning(
        "slide %d 圖表渲染失敗，fallback 為文字顯示 mermaid 原始碼。",
        slide_data.slide_number,
    )

    fallback_bullets = ["[圖表渲染失敗，mermaid 原始碼如下]"] + mermaid_string.split("\n")
    left = int(MARGIN_LEFT)
    top = int(MARGIN_TOP_CONTENT)
    width = slide_width - int(MARGIN_LEFT) - int(MARGIN_RIGHT)
    height = slide_height - top - int(MARGIN_RIGHT)

    _add_bullet_textbox(slide, fallback_bullets, left, top, width, height)


def _handle_bullets_with_diagram(
    slide,
    slide_data: Slide,
    slide_width: int,
    slide_height: int,
) -> None:
    """處理 bullets_with_diagram：左側條列 + 右側 PNG 圖表。

    若圖表渲染失敗，右側改為顯示 mermaid 原始碼的文字框。

    Args:
        slide: pptx Slide 物件
        slide_data: Slide 資料
        slide_width: 投影片寬度（EMU）
        slide_height: 投影片高度（EMU）
    """
    bullets = slide_data.content.bullets or []
    mermaid_string = slide_data.content.mermaid or ""

    top = int(MARGIN_TOP_CONTENT)
    content_height = slide_height - top - int(MARGIN_RIGHT)

    # 左側條列（45% 寬）
    left_left = int(MARGIN_LEFT)
    left_width = int(slide_width * 0.45)

    _add_bullet_textbox(
        slide=slide,
        bullets=bullets,
        left=left_left,
        top=top,
        width=left_width,
        height=content_height,
    )

    # 右側圖表（50% 寬，從 50% 位置開始）
    right_left = int(slide_width * 0.50)
    right_width = int(slide_width * 0.50) - int(MARGIN_RIGHT)

    output_stem = f"slide_{slide_data.slide_number}_diagram"
    result = build_diagram(
        mermaid_string=mermaid_string,
        output_stem=output_stem,
        output_dir=DIAGRAM_OUTPUT_DIR,
        config_path=DIAGRAM_CONFIG_PATH,
    )

    if result is not None:
        slide.shapes.add_picture(
            str(result.png_path),
            right_left,
            top,
            width=right_width,
        )
        return

    # 圖表渲染失敗：右側顯示 mermaid 原始碼
    logger.warning(
        "slide %d 圖表渲染失敗，fallback 為文字顯示 mermaid 原始碼。",
        slide_data.slide_number,
    )

    fallback_bullets = ["[圖表渲染失敗]"] + mermaid_string.split("\n")

    _add_bullet_textbox(
        slide=slide,
        bullets=fallback_bullets,
        left=right_left,
        top=top,
        width=right_width,
        height=content_height,
    )


# ============================================================
# 文字框與表格 helper
# ============================================================


def _add_bullet_textbox(
    slide,
    bullets: list[str],
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    """在投影片上新增條列文字框。

    Args:
        slide: pptx Slide 物件
        bullets: 條列文字 list
        left: 左邊緣位置（EMU）
        top: 上邊緣位置（EMU）
        width: 文字框寬度（EMU）
        height: 文字框高度（EMU）
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        # 第一段直接用 text_frame 的第一個 paragraph
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        run = para.add_run()
        run.text = f"{BULLET_CHAR}  {bullet_text}"
        run.font.name = FONT_BODY
        run.font.size = FONT_SIZE_BODY
        run.font.color.rgb = COLOR_BLACK


def _add_styled_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: int,
    top: int,
    width: int,
    row_count: int,
    col_count: int,
    slide_height: int,
) -> None:
    """在投影片上新增帶樣式的表格。

    標題列：公司主色背景（#005BAC）、白色粗體文字。
    資料列：交替底色（白 / 淺灰）、深色文字。

    Args:
        slide: pptx Slide 物件
        headers: 欄位名稱 list
        rows: 資料列的二維 list
        left: 左邊緣位置（EMU）
        top: 上邊緣位置（EMU）
        width: 表格寬度（EMU）
        row_count: 總列數（含標題列）
        col_count: 欄數
        slide_height: 投影片高度（EMU），用於計算表格可用高度
    """
    # 估算列高：平均分配表格高度（留 55% 空間給表格）
    available_height = int(slide_height * 0.55)
    row_height = available_height // max(row_count, 1)

    table = slide.shapes.add_table(
        rows=row_count,
        cols=col_count,
        left=left,
        top=top,
        width=width,
        height=available_height,
    ).table

    # 均分欄寬
    col_width = width // max(col_count, 1)

    for col_idx in range(col_count):
        table.columns[col_idx].width = col_width

    # 填入標題列
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        _style_table_cell(
            cell=cell,
            text=header_text,
            is_header=True,
            is_even_row=False,
        )

    # 填入資料列
    for row_idx, row_data in enumerate(rows):
        # row_idx 從 0 開始，表格中列 index 從 1 開始（0 是標題列）
        table_row_idx = row_idx + 1
        is_even = (row_idx % 2 == 0)

        for col_idx in range(col_count):
            cell = table.cell(table_row_idx, col_idx)

            # 若資料列欄位數不足，填空字串
            if col_idx < len(row_data):
                cell_text = row_data[col_idx]
            else:
                cell_text = ""

            _style_table_cell(
                cell=cell,
                text=cell_text,
                is_header=False,
                is_even_row=is_even,
            )


def _style_table_cell(
    cell,
    text: str,
    is_header: bool,
    is_even_row: bool,
) -> None:
    """設定表格 cell 的樣式（背景色、字型、粗體、對齊）。

    Args:
        cell: pptx TableCell 物件
        text: 填入的文字內容
        is_header: 是否為標題列
        is_even_row: 是否為偶數資料列（用於交替底色）
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    # 設定背景色
    if is_header:
        bg_color = COLOR_PRIMARY
    elif is_even_row:
        bg_color = COLOR_ROW_EVEN
    else:
        bg_color = COLOR_ROW_ODD

    # 透過 XML 設定 cell 背景色（python-pptx 無直接 API）
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # 移除舊的 solidFill（避免重複設定）
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)

    solid_fill = etree.SubElement(tcPr, qn("a:solidFill"))
    srg_clr = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srg_clr.set("val", str(bg_color))

    # 設定文字內容與字型
    tf = cell.text_frame
    tf.word_wrap = True

    # 清除現有段落內容
    para = tf.paragraphs[0]
    para.clear()

    run = para.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = FONT_SIZE_TITLE_ROW if is_header else FONT_SIZE_BODY
    run.font.bold = is_header

    # 文字顏色
    if is_header:
        run.font.color.rgb = COLOR_WHITE
    else:
        run.font.color.rgb = COLOR_BLACK

    # 對齊
    para.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
