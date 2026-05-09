"""
掃描 PPTX 模板的 placeholder 結構，供 build_pptx.py 開發參考。

新增功能：自動生成 template_map.json 初稿（使用者再手動調整 layout_type 對應）。

使用方式：
    python scan_template.py
    python scan_template.py --dump-map        # 額外輸出 template_map.json 初稿
    python scan_template.py --dump-map --out config/template_map_new.json
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


# 模板路徑（相對於專案根目錄）
DEFAULT_TEMPLATE = Path("template") / "Creative-Idea-Bulb-PowerPoint-Template.pptx"

# template_map.json 輸出預設路徑
DEFAULT_MAP_OUTPUT = Path("config") / "template_map_draft.json"

# 本專案使用的 layout_type（需要手動對應到 layout_name）
LAYOUT_TYPES = [
    "title_slide",
    "section_divider",
    "bullets_only",
    "table_only",
    "bullets_with_table",
    "diagram_only",
    "bullets_with_diagram",
]


def scan_template(template_path: Path) -> None:
    """掃描並列印模板的 layout 與 placeholder 資訊。

    Args:
        template_path: PPTX 模板的 Path
    """
    prs = Presentation(str(template_path))

    # 列印投影片尺寸
    print("=== 投影片尺寸 ===")
    print(f"寬度: {prs.slide_width} EMU = {prs.slide_width / 914400:.2f} inches")
    print(f"高度: {prs.slide_height} EMU = {prs.slide_height / 914400:.2f} inches")
    print()

    # 列印 slide_layouts（prs.slide_layouts 中的官方清單）
    print("=== Slide Layouts（prs.slide_layouts）===")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  Layout[{i}]: {layout.name}")
        for ph in layout.placeholders:
            fmt = ph.placeholder_format
            print(
                f"    idx={fmt.idx}, type={fmt.type}, name=\"{ph.name}\", "
                f"size=({ph.width},{ph.height}), pos=({ph.left},{ph.top})"
            )
        print()

    # 列印現有 slides 並收集額外的 layout
    print("=== 現有 Slides ===")
    extra_layouts: dict = {}

    for j, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"  Slide[{j}]: layout=\"{layout_name}\"")

        for ph in slide.placeholders:
            text_preview = ph.text[:80] if ph.text else "(empty)"
            print(f"    idx={ph.placeholder_format.idx}, name=\"{ph.name}\", text=\"{text_preview}\"")

        print()

        # 收集不在 slide_layouts 清單中的額外 layout
        known_names = {l.name for l in prs.slide_layouts}
        if layout_name not in known_names and layout_name not in extra_layouts:
            extra_layouts[layout_name] = slide.slide_layout

    # 若有額外 layout，也列印其 placeholder 資訊
    if extra_layouts:
        print("=== 額外 Layouts（僅存在於現有 Slides 中，不在 slide_layouts 清單）===")
        for name, layout in extra_layouts.items():
            print(f"  Layout: {name}")
            for ph in layout.placeholders:
                fmt = ph.placeholder_format
                print(
                    f"    idx={fmt.idx}, type={fmt.type}, name=\"{ph.name}\", "
                    f"size=({ph.width},{ph.height}), pos=({ph.left},{ph.top})"
                )
            print()


def generate_template_map_draft(
    template_path: Path,
    output_path: Path,
) -> None:
    """生成 template_map.json 初稿。

    從模板中收集所有可用的 layout_name（含 slide_layouts 和現有 slides 中的 layout），
    並以預設值填入所有 layout_type 的對應設定。

    使用者需手動調整各 layout_type 對應到正確的 layout_name。

    Args:
        template_path: PPTX 模板的 Path
        output_path: 輸出的 template_map.json 路徑
    """
    prs = Presentation(str(template_path))

    # 收集所有可用的 layout_name（含官方清單和現有 slides 中的 layout）
    all_layout_names: list[str] = []

    for layout in prs.slide_layouts:
        if layout.name not in all_layout_names:
            all_layout_names.append(layout.name)

    for slide in prs.slides:
        name = slide.slide_layout.name
        if name not in all_layout_names:
            all_layout_names.append(name)

    # 建立初稿 template_map 結構
    # 預設使用第一個可用的 layout_name 和 idx=10, 11
    default_layout_name = all_layout_names[0] if all_layout_names else "Cover Slide layout"

    layouts_draft: dict = {}

    for lt in LAYOUT_TYPES:
        layouts_draft[lt] = {
            "layout_name": default_layout_name,
            "title_idx": 10,
            "subtitle_idx": 11,
        }

    template_map_draft = {
        "template_path": str(template_path).replace("\\", "/"),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "_available_layouts": all_layout_names,
        "_note": "請將 layouts 中各 layout_type 的 layout_name 改為對應的模板 layout 名稱。",
        "layouts": layouts_draft,
    }

    output_path.parent.mkdir(parents=True, exist_ok=True)

    with output_path.open("w", encoding="utf-8") as f:
        json.dump(template_map_draft, f, ensure_ascii=False, indent=4)

    print(f"template_map.json 初稿已輸出至：{output_path}")
    print(f"可用的 layout 清單：{all_layout_names}")
    print("請手動調整各 layout_type 的 layout_name 對應。")


def main() -> None:
    """主程式入口：解析參數並執行掃描。"""
    parser = argparse.ArgumentParser(
        description="掃描 PPTX 模板的 placeholder 結構"
    )

    parser.add_argument(
        "--template",
        type=Path,
        default=DEFAULT_TEMPLATE,
        help=f"PPTX 模板路徑（預設：{DEFAULT_TEMPLATE}）",
    )

    parser.add_argument(
        "--dump-map",
        action="store_true",
        help="額外生成 template_map.json 初稿",
    )

    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_MAP_OUTPUT,
        help=f"template_map.json 初稿的輸出路徑（預設：{DEFAULT_MAP_OUTPUT}）",
    )

    args = parser.parse_args()

    if not args.template.exists():
        print(f"錯誤：模板檔案不存在：{args.template}")
        return

    scan_template(args.template)

    if args.dump_map:
        print()
        print("=== 生成 template_map.json 初稿 ===")
        generate_template_map_draft(args.template, args.out)


if __name__ == "__main__":
    main()
